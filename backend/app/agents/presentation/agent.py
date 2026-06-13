"""
Agent 10: Presentation Agent
Generates a PPTX from research findings.

Fixed: _fill_background used broken slide.part.package.presentation API.
Now uses prs.slide_width / prs.slide_height directly (passed in).
"""
import json
from pathlib import Path
from typing import List, Dict, Any

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu
from loguru import logger

from app.utils.gemini_client import ask_gemini
from app.core.config import settings

# ── Palette ───────────────────────────────────────────────────────────────────
C_BG      = RGBColor(0x0F, 0x17, 0x2A)  # deep navy
C_HEADER  = RGBColor(0x1E, 0x29, 0x3B)  # slate-800
C_ACCENT  = RGBColor(0x38, 0x7E, 0xD4)  # blue
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY    = RGBColor(0x94, 0xA3, 0xB8)  # slate-400
C_LIGHT   = RGBColor(0xCB, 0xD5, 0xE1)  # slate-300


class PresentationAgent:
    def __init__(self):
        self.storage_dir = Path(settings.PDF_STORAGE_DIR).parent / "presentations"
        self.storage_dir.mkdir(parents=True, exist_ok=True)

    async def run(
        self,
        topic: str,
        project_id: int,
        papers: List[Dict[str, Any]],
        trends: str = "",
        gaps: str = "",
        comparison: str = "",
    ) -> Dict[str, Any]:
        slide_data = await self._generate_slide_content(topic, papers, trends, gaps)
        file_path  = self._build_pptx(project_id, slide_data, papers)
        logger.info(f"[PresentationAgent] Saved: {file_path}")
        return {"file_path": str(file_path), "slide_data": slide_data}

    # ── AI slide content ──────────────────────────────────────────────────────

    async def _generate_slide_content(
        self, topic: str, papers: List[Dict], trends: str, gaps: str
    ) -> Dict:
        paper_list = "\n".join(
            f"- {p.get('title','')[:90]} ({p.get('year','')})"
            for p in papers[:10]
        )
        prompt = f"""Generate a research presentation outline for "{topic}".

Papers:
{paper_list}

Trends: {trends[:400] if trends else "N/A"}
Gaps: {gaps[:400] if gaps else "N/A"}

Return ONLY valid JSON (no markdown fences):
{{
  "title": "short presentation title",
  "subtitle": "subtitle line",
  "slides": [
    {{"title": "Agenda", "bullets": ["Introduction", "Papers Reviewed", "Key Findings", "Trends & Gaps", "Conclusion"]}},
    {{"title": "Introduction", "bullets": ["3-4 concise bullet points about the topic"]}},
    {{"title": "Papers Reviewed", "bullets": ["Author et al. (year) — one line summary per paper, up to 8"]}},
    {{"title": "Key Findings", "bullets": ["3-5 most important findings"]}},
    {{"title": "Research Trends", "bullets": ["3-4 trend bullets"]}},
    {{"title": "Research Gaps", "bullets": ["3-4 gap bullets"]}},
    {{"title": "Conclusion", "bullets": ["2-3 takeaway bullets"]}}
  ]
}}"""

        response = await ask_gemini(prompt, max_tokens=2000)
        try:
            clean = response.strip()
            # Strip any markdown fences robustly
            if "```" in clean:
                parts = clean.split("```")
                for part in parts:
                    p = part.strip().lstrip("json").strip()
                    if p.startswith("{"):
                        clean = p
                        break
            return json.loads(clean)
        except Exception as e:
            logger.warning(f"[PresentationAgent] JSON parse failed: {e}")
            return self._fallback_slides(topic, papers)

    # ── PPTX builder ──────────────────────────────────────────────────────────

    def _build_pptx(self, project_id: int, slide_data: Dict, papers: List[Dict]) -> Path:
        prs = PptxPresentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Store dimensions for helper methods
        self._w = prs.slide_width
        self._h = prs.slide_height

        blank = prs.slide_layouts[6]  # blank layout

        self._title_slide(prs, blank, slide_data)
        for s in slide_data.get("slides", []):
            self._content_slide(prs, blank, s)
        self._papers_slide(prs, blank, papers)

        out = self.storage_dir / f"project_{project_id}_presentation.pptx"
        prs.save(str(out))
        return out

    def _fill_bg(self, slide, color: RGBColor):
        """Add a full-slide background rectangle as the bottom-most shape."""
        shape = slide.shapes.add_shape(1, Emu(0), Emu(0), self._w, self._h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        # Send to back by moving element to first position in spTree
        sp_tree = slide.shapes._spTree
        sp_tree.remove(shape._element)
        sp_tree.insert(2, shape._element)   # index 2 = after spTree bookmarks

    def _title_slide(self, prs, layout, slide_data: Dict):
        slide = prs.slides.add_slide(layout)
        self._fill_bg(slide, C_BG)

        # Accent bar across middle
        bar = slide.shapes.add_shape(1, Inches(0), Inches(3.35), self._w, Inches(0.07))
        bar.fill.solid()
        bar.fill.fore_color.rgb = C_ACCENT
        bar.line.fill.background()

        # Title
        tb = slide.shapes.add_textbox(Inches(1.2), Inches(1.4), Inches(10.9), Inches(1.7))
        tf = tb.text_frame; tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.text           = slide_data.get("title", "Research Review")
        p.font.size      = Pt(38)
        p.font.bold      = True
        p.font.color.rgb = C_WHITE
        p.alignment      = PP_ALIGN.CENTER

        # Subtitle
        tb2 = slide.shapes.add_textbox(Inches(1.2), Inches(3.55), Inches(10.9), Inches(0.9))
        tf2 = tb2.text_frame
        p2  = tf2.paragraphs[0]
        p2.text           = slide_data.get("subtitle", "AI-Powered Literature Review")
        p2.font.size      = Pt(20)
        p2.font.color.rgb = C_GRAY
        p2.alignment      = PP_ALIGN.CENTER

    def _content_slide(self, prs, layout, info: Dict):
        slide = prs.slides.add_slide(layout)
        self._fill_bg(slide, C_BG)

        # Header band
        hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), self._w, Inches(1.15))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = C_HEADER
        hdr.line.fill.background()

        # Left accent stripe
        stripe = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.07), Inches(1.15))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = C_ACCENT
        stripe.line.fill.background()

        # Slide title
        tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.18), Inches(12.5), Inches(0.8))
        tf = tb.text_frame
        p  = tf.paragraphs[0]
        p.text           = info.get("title", "")
        p.font.size      = Pt(26)
        p.font.bold      = True
        p.font.color.rgb = C_WHITE

        # Bullet content
        ct = slide.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(12.1), Inches(5.9))
        tf2 = ct.text_frame; tf2.word_wrap = True
        bullets = info.get("bullets", [])
        for i, bullet in enumerate(bullets):
            para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            # Bullet dot + text
            run = para.add_run()
            run.text           = f"  {bullet}"
            run.font.size      = Pt(17)
            run.font.color.rgb = C_LIGHT
            para.space_before  = Pt(7)

    def _papers_slide(self, prs, layout, papers: List[Dict]):
        slide = prs.slides.add_slide(layout)
        self._fill_bg(slide, C_BG)

        hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), self._w, Inches(1.15))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = C_HEADER
        hdr.line.fill.background()

        stripe = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.07), Inches(1.15))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = C_ACCENT
        stripe.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.18), Inches(12.5), Inches(0.8))
        p  = tb.text_frame.paragraphs[0]
        p.text = "Papers Analyzed"
        p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = C_WHITE

        ct = slide.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(12.1), Inches(5.9))
        tf2 = ct.text_frame; tf2.word_wrap = True
        for i, paper in enumerate(papers[:9]):
            para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            run  = para.add_run()
            title = (paper.get("title") or "Unknown")[:75]
            year  = paper.get("year", "")
            run.text           = f"  {title}  ({year})"
            run.font.size      = Pt(14)
            run.font.color.rgb = C_GRAY
            para.space_before  = Pt(5)

    def _fallback_slides(self, topic: str, papers: List[Dict]) -> Dict:
        return {
            "title":    f"Literature Review: {topic}",
            "subtitle": "AI-Generated Research Analysis",
            "slides": [
                {"title": "Overview",   "bullets": [f"Topic: {topic}", f"{len(papers)} papers analyzed"]},
                {"title": "Papers",     "bullets": [p.get("title","")[:70] for p in papers[:6]]},
                {"title": "Conclusion", "bullets": ["See full literature review for details"]},
            ],
        }
